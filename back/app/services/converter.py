import os
import subprocess
import uuid
import io
from app.core.config import settings
from fastapi import UploadFile
from app.schemas.document import DocumentCreate, DocumentRead
from app.crud.document import create_document, get_document
from sqlalchemy.ext.asyncio import AsyncSession
from openpyxl.workbook import Workbook
from openpyxl.drawing.image import Image as XlImage
from app.core.exceptions.global_exception import GlobalException
from app.core.exceptions.error_code import ErrorCode
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path


PATH = os.path.abspath(settings.FILE_STORAGE_PATH)
LIBREOFFICE_PATH = settings.LIBREOFFICE_PATH
POPPLER_PATH = settings.POPPLER_PATH

async def convert_document(input_ext: str, output_ext: str, file: UploadFile, db: AsyncSession) -> DocumentRead:
    # 디렉토리가 없으면 새로 생성
    os.makedirs(PATH, exist_ok=True)

    input_filename = file.filename
    output_filename = input_filename.replace(f".{input_ext}", f".{output_ext}")

    # 경로 + uuid 생성
    save_input_filepath = os.path.join(PATH, get_uuid_filename(input_ext))
    convert_filepath = save_input_filepath.replace(f".{input_ext}", f".{output_ext}")

    # DTO 생성 is_success를 기본 False로
    doc = DocumentCreate(output_filename=output_filename,
                          convert_filename=os.path.basename(convert_filepath), 
                          is_zip=False, is_success=False)
    
    # 요청 확장자와 실제 파일의 확장자가 다르면 is_success를 False로 db 저장 후 반환
    if f".{input_ext}" != os.path.splitext(input_filename)[1]:
       return await create_document(create=doc, db=db)

    try:
      # 파일 업로드
      with open(save_input_filepath, "wb") as f:
        f.write(await file.read())
      
      # 변환
      if input_ext == "pdf":
        await from_pdf(save_input_filepath=save_input_filepath, 
                  convert_filepath=convert_filepath, output_ext=output_ext)
      else:
        await to_pdf(output_ext=output_ext, save_input_filepath=save_input_filepath)
      
      # 성공 시 success를 True로
      doc.is_success = True
    except Exception as e:
      print(e)
    finally:
      # 성공, 실패 상관없이 DB 저장
      return await create_document(create=doc, db=db)


async def docs_download(convert_filename: str, db: AsyncSession):
  doc = await get_document(convert_filename=convert_filename, db=db)

  if not doc:
    raise GlobalException(ErrorCode.NOT_FOUND_DOCUMENTS)

  if not doc.is_success:
    raise GlobalException(ErrorCode.CONVERSION_FAILED_DOCUMENTS)
  
  convert_filepath = os.path.join(PATH, doc.convert_filename)

  if not os.path.exists(convert_filepath):
    raise GlobalException(ErrorCode.NOT_FOUND_DOCUMENTS)

  return doc.output_filename, convert_filepath


async def from_pdf(save_input_filepath: str, convert_filepath: str, output_ext: str):
  images = convert_from_path(save_input_filepath, dpi=200, poppler_path=POPPLER_PATH)
  
  if output_ext == "docx":
    doc = Document()

    for i, image in enumerate(images):
        # 메모리에서 이미지 처리
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        
        # 이미지를 문서에 추가 (페이지 너비에 맞게)
        doc.add_picture(img_byte_arr, width=Inches(6))
        
        # 마지막 페이지가 아니면 페이지 나누기 추가
        if i < len(images) - 1:
            doc.add_page_break()

    doc.save(convert_filepath)
  
  if output_ext == "pptx":
    print("podf -> pptx 리팩토링")
    # 새 PowerPoint 프레젠테이션 생성
    prs = Presentation()

    # 첫 페이지 이미지 기준으로 슬라이드 크기 설정
    first_image = images[0]
    img_width, img_height = first_image.size  # pixel 단위
    prs.slide_width = Inches(img_width / 200 * 72 / 96)   # pixel → point → inch
    prs.slide_height = Inches(img_height / 200 * 72 / 96)
    
    # 각 이미지를 새 슬라이드에 추가
    for i, image in enumerate(images):
        # 빈 슬라이드 생성
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6은 빈 레이아웃
        
        # 이미지를 메모리 스트림으로 변환
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        
        # 슬라이드 크기 계산
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 이미지를 슬라이드에 추가 (슬라이드 전체 크기에 맞게)
        slide.shapes.add_picture(
            img_byte_arr,
            left=0,
            top=0,
            width=slide_width,
            height=slide_height
        )
    
    # 프레젠테이션 저장
    prs.save(convert_filepath)

  if output_ext == "xlsx":
    # 새 Excel 워크북 생성
    wb = Workbook()
    ws = wb.active
    
    # 각 이미지를 동일한 시트에 세로로 배치
    current_row = 1
    
    for image in images:
        # 이미지 크기 조절
        img_width, img_height = image.size
        scaling_factor = min(800 / img_width, 500 / img_height)
        new_width = int(img_width * scaling_factor)
        new_height = int(img_height * scaling_factor)
        resized_image = image.resize((new_width, new_height))

        # 이미지 메모리 스트림 변환
        img_byte_arr = io.BytesIO()
        resized_image.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)

        # 엑셀 이미지 객체 생성 및 삽입
        img = XlImage(img_byte_arr)
        cell_position = f"A{current_row}"
        ws.add_image(img, cell_position)

        # 행 높이 조정 (이미지 높이 → 엑셀 행 수로 변환)
        # openpyxl 기준: 약 1행당 0.75pt ≒ 1px 정도로 보면 무난
        px_per_row = 15  # 15px 정도를 한 행에 배치한다고 가정
        image_height_in_rows = max(1, int(new_height / px_per_row))  # 최소 1행

        # 다음 이미지 배치할 행 계산
        for row in range(current_row, current_row + image_height_in_rows):
            ws.row_dimensions[row].height = 15  # 기본 행 높이 설정

        current_row += image_height_in_rows + 2  # 이미지 + 간격
    
    # 워크북 저장
    wb.save(convert_filepath)


async def to_pdf(output_ext: str, save_input_filepath: str):
  subprocess.run([
     LIBREOFFICE_PATH,
     "--headless",
     "--convert-to", output_ext,
     "--outdir", PATH,
     save_input_filepath
     ])

def get_uuid_filename(ext: str):
   return f"{uuid.uuid4().hex}.{ext}"