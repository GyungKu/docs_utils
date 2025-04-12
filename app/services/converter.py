import os
import subprocess
import uuid
import pdfplumber
from app.core.config import settings
from fastapi import UploadFile
from app.schemas.document import DocumentCreate, DocumentRead
from app.crud.document import create_document, get_document
from sqlalchemy.ext.asyncio import AsyncSession
from pdf2docx import Converter
from pptx import Presentation
from pdf2image import convert_from_path
from pptx.util import Inches
from openpyxl.workbook import Workbook
from app.core.exceptions.global_exception import GlobalException
from app.core.exceptions.error_code import ErrorCode


PATH = os.path.abspath(settings.FILE_STORAGE_PATH)
LIBREOFFICE_PATH = settings.LIBREOFFICE_PATH

async def convert_document(input_ext: str, output_ext: str, file: UploadFile, db: AsyncSession) -> DocumentRead:
    # 디렉토리가 없으면 새로 생성
    os.makedirs(PATH, exist_ok=True)

    input_filename = file.filename
    output_filename = input_filename.replace(f".{input_ext}", f".{output_ext}")

    # 경로 + uuid 생성
    save_input_filepath = os.path.join(PATH, get_uuid_filename(input_ext))
    convert_filepath = save_input_filepath.replace(f".{input_ext}", f".{output_ext}")

    # DTO 생성 is_success를 기본 False로
    doc = DocumentCreate(output_filename=output_filename,
                          convert_filename=os.path.basename(convert_filepath), 
                          is_zip=False, is_success=False)
    
    # 요청 확장자와 실제 파일의 확장자가 다르면 is_success를 False로 db 저장 후 반환
    if f".{input_ext}" != os.path.splitext(input_filename)[1]:
       return await create_document(create=doc, db=db)

    try:
      # 파일 업로드
      with open(save_input_filepath, "wb") as f:
        f.write(await file.read())
      
      # 변환
      if input_ext == "pdf":
        await from_pdf(save_input_filepath=save_input_filepath, 
                  convert_filepath=convert_filepath, output_ext=output_ext)
      else:
        await to_pdf(output_ext=output_ext, save_input_filepath=save_input_filepath)
      
      # 성공 시 success를 True로
      doc.is_success = True
    finally:
      # 성공, 실패 상관없이 DB 저장
      return await create_document(create=doc, db=db)


async def docs_download(id: int, db: AsyncSession):
  doc = await get_document(document_id=id, db=db)

  if not doc:
    raise GlobalException(ErrorCode.NOT_FOUND_DOCUMENTS)

  if not doc.is_success:
    raise GlobalException(ErrorCode.CONVERSION_FAILED_DOCUMENTS)
  
  convert_filepath = os.path.join(PATH, doc.convert_filename)

  if not os.path.exists(convert_filepath):
    raise GlobalException(ErrorCode.NOT_FOUND_DOCUMENTS)

  return doc.output_filename, convert_filepath


async def from_pdf(save_input_filepath: str, convert_filepath: str, output_ext: str):
  if output_ext == "docx":
    cv = Converter(save_input_filepath)
    cv.convert(convert_filepath, start=0, end=None)
    cv.close()
  
  if output_ext == "pptx":
    images = convert_from_path(save_input_filepath, dpi=200)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for image in images:
      slide = prs.slides.add_slide(blank_slide_layout)
      img_path = os.path.join(PATH, get_uuid_filename("jpg"))
      image.save(img_path, "JPEG")
      slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                               width=prs.slide_width, height=prs.slide_height)
      os.remove(img_path)

    prs.save(convert_filepath)

  if output_ext == "xlsx":
    wb = Workbook()
    ws = wb.active

    with pdfplumber.open(save_input_filepath) as pdf:
      for page in pdf.pages:
        tables = page.extract_tables()
        for table in tables:
          for row in table:
            ws.append(row)

    wb.save(convert_filepath)


async def to_pdf(output_ext: str, save_input_filepath: str):
  subprocess.run([
     LIBREOFFICE_PATH,
     "--headless",
     "--convert-to", output_ext,
     "--outdir", PATH,
     save_input_filepath
     ])

def get_uuid_filename(ext: str):
   return f"{uuid.uuid4().hex}.{ext}"